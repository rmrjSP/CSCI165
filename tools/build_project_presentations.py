from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\CSCI165")

BG = RGBColor(248, 249, 251)
NAVY = RGBColor(27, 43, 65)
BLUE = RGBColor(59, 93, 201)
TEAL = RGBColor(34, 139, 154)
GRAY = RGBColor(92, 99, 112)
LIGHT = RGBColor(228, 233, 242)
WHITE = RGBColor(255, 255, 255)


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def set_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_header(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(11.9), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = NAVY

    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.58), Inches(0.95), Inches(11.5), Inches(0.35))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = GRAY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.22), Inches(2.3), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_title_slide(prs: Presentation, title: str, subtitle: str, course_line: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.9), Inches(11.2), Inches(4.7))
    banner.fill.solid()
    banner.fill.fore_color.rgb = WHITE
    banner.line.color.rgb = LIGHT

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(0.9), Inches(0.22), Inches(4.7))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.25), Inches(1.45), Inches(9.8), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY

    sub_box = slide.shapes.add_textbox(Inches(1.28), Inches(3.0), Inches(8.8), Inches(0.8))
    p2 = sub_box.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = "Aptos"
    r2.font.size = Pt(18)
    r2.font.color.rgb = TEAL

    course_box = slide.shapes.add_textbox(Inches(1.28), Inches(4.1), Inches(8.8), Inches(0.8))
    p3 = course_box.text_frame.paragraphs[0]
    r3 = p3.add_run()
    r3.text = course_line
    r3.font.name = "Aptos"
    r3.font.size = Pt(14)
    r3.font.color.rgb = GRAY


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title, subtitle)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.55), Inches(11.0), Inches(5.45))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT

    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(10.3), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"• {bullet}"
        p.font.name = "Aptos"
        p.font.size = Pt(20)
        p.font.color.rgb = NAVY
        p.space_after = Pt(10)


def add_two_column_slide(prs: Presentation, title: str, left_title: str, left_bullets: list[str], right_title: str, right_bullets: list[str], subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title, subtitle)

    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.55), Inches(5.35), Inches(5.45))
    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.1), Inches(1.55), Inches(5.25), Inches(5.45))
    for shape in (left, right):
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LIGHT

    def fill_column(x, y, w, h, header, bullets):
        hb = slide.shapes.add_textbox(x, y, w, Inches(0.45))
        p = hb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = header
        r.font.name = "Aptos Display"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = BLUE

        body = slide.shapes.add_textbox(x, y + Inches(0.5), w, h - Inches(0.6))
        tf = body.text_frame
        tf.word_wrap = True
        first = True
        for bullet in bullets:
            pp = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            pp.text = f"• {bullet}"
            pp.font.name = "Aptos"
            pp.font.size = Pt(17)
            pp.font.color.rgb = NAVY
            pp.space_after = Pt(8)

    fill_column(Inches(0.95), Inches(1.85), Inches(4.7), Inches(4.9), left_title, left_bullets)
    fill_column(Inches(6.4), Inches(1.85), Inches(4.55), Inches(4.9), right_title, right_bullets)


def add_image_slide(prs: Presentation, title: str, image_paths: list[Path], captions: list[str], subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title, subtitle)

    if len(image_paths) == 1:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.55), Inches(10.9), Inches(5.45))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT
        slide.shapes.add_picture(str(image_paths[0]), Inches(1.0), Inches(1.85), width=Inches(10.4), height=Inches(4.45))
        cap = slide.shapes.add_textbox(Inches(1.0), Inches(6.4), Inches(10.2), Inches(0.35))
        p = cap.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = captions[0]
        r.font.name = "Aptos"
        r.font.size = Pt(11)
        r.font.color.rgb = GRAY
        return

    coords = [(Inches(0.7), Inches(1.75)), (Inches(6.2), Inches(1.75))]
    for idx, (img, cap_text) in enumerate(zip(image_paths, captions)):
        x, y = coords[idx]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y - Inches(0.15), Inches(5.0), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT
        slide.shapes.add_picture(str(img), x + Inches(0.15), y, width=Inches(4.7), height=Inches(4.05))
        cap = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(4.18), Inches(4.8), Inches(0.45))
        p = cap.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = cap_text
        r.font.name = "Aptos"
        r.font.size = Pt(10.5)
        r.font.color.rgb = GRAY


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title, subtitle)

    cols = len(headers)
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.65), Inches(1.65), Inches(11.0), Inches(5.1)).table

    for i, head in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = WHITE

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else RGBColor(242, 245, 250)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(11)
                run.font.color.rgb = NAVY


def build_8queens():
    project = ROOT / "8Queens"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    summary = read_csv(project / "results" / "summary.csv")

    add_title_slide(
        prs,
        "Incremental Evolutionary Algorithm for 8-Queens",
        "Motivation, method, evaluation, and results",
        "CSCI 165 project presentation"
    )

    add_two_column_slide(
        prs,
        "Motivation And Problem Statement",
        "Why This Domain Matters",
        [
            "8-Queens is a classic constraint satisfaction and search problem.",
            "It is small enough to visualize clearly but still rich enough to compare optimization strategies.",
            "It highlights how representation choices can shrink a search space dramatically."
        ],
        "Problem Addressed",
        [
            "Find a placement of 8 queens on a chessboard so that no two attack each other.",
            "Compare how an evolutionary algorithm behaves under different population, mutation, and tournament settings.",
            "Measure success rate, solution quality, and generations needed to solve."
        ]
    )

    add_two_column_slide(
        prs,
        "Background And Approach",
        "Background",
        [
            "A board is valid when no pair of queens shares a row or diagonal.",
            "The project uses a permutation encoding, so each column has exactly one queen and rows never repeat.",
            "Fitness = 28 minus the number of attacking pairs, so 28 is a perfect solution."
        ],
        "Method",
        [
            "Initialize a population of random row permutations.",
            "Use tournament selection, cut-and-crossfill crossover, and swap mutation.",
            "Carry forward elite individuals and repeat until a perfect board is found or the generation budget ends."
        ]
    )

    add_bullets_slide(
        prs,
        "Experiment Setup",
        [
            "30 runs per configuration with max_gens = 1000 and elitism = 2.",
            "Population sweep: 20, 50, 100 with mutation = 0.10 and tournament = 3.",
            "Mutation sweep: 0.05, 0.10, 0.20 with population = 100 and tournament = 3.",
            "Tournament sweep: 2, 3, 5 with population = 100 and mutation = 0.10.",
            "Metrics: success rate, mean best fitness, average generation solved, and runtime."
        ]
    )

    best_rows = [
        [r["config"], r["success_rate"], r["mean_fitness"], r["avg_gen_solved"], r["mean_time_ms"]]
        for r in summary
    ]
    add_table_slide(
        prs,
        "Evaluation Results",
        ["Config", "Success", "Mean Fitness", "Avg Gen Solved", "Time (ms)"],
        best_rows,
        "Summary across 30 runs per configuration"
    )

    add_image_slide(
        prs,
        "Visual Results",
        [project / "figures" / "success_rates.png", project / "figures" / "convergence_plot.png"],
        [
            "Success rates show `pop=20` was the only configuration with noticeable failures.",
            "Convergence curves show high-performing settings reaching perfect fitness quickly."
        ]
    )

    add_image_slide(
        prs,
        "Representation And Example Boards",
        [project / "figures" / "search_space.png", project / "figures" / "example_boards.png"],
        [
            "Permutation encoding reduces the search space from billions of boards to 40,320 candidates.",
            "The board visualization makes the conflict structure and solved state easy to explain."
        ]
    )

    add_bullets_slide(
        prs,
        "Conclusions And Future Work",
        [
            "The incremental permutation encoding was the key design choice because it reduced wasted search substantially.",
            "Most tested settings solved the problem perfectly in all 30 runs; `pop=20` was the weakest configuration at 23/30 successes.",
            "Population sizes of 50 and 100 gave perfect success with fast convergence, while very large tournament pressure slowed solving.",
            "Future work: test larger N-Queens variants, compare against hill climbing or simulated annealing, and study sensitivity to crossover and elitism."
        ]
    )

    out = project / "8Queens_project_presentation.pptx"
    prs.save(out)
    return out


def build_rastrigin():
    project = ROOT / "Rastrigin"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    summary = read_csv(project / "results" / "summary.csv")
    rows = {r["algorithm"]: r for r in summary}

    add_title_slide(
        prs,
        "Gradient Descent vs Simulated Annealing on Rastrigin",
        "Comparing local optimization and stochastic search on a multimodal landscape",
        "CSCI 165 project presentation"
    )

    add_two_column_slide(
        prs,
        "Motivation And Problem Statement",
        "Why This Domain Matters",
        [
            "Optimization methods are used in machine learning, engineering design, scheduling, and many search tasks.",
            "Rastrigin is a standard benchmark because it has many local minima and a known global optimum.",
            "It is a good test of whether an algorithm exploits quickly or explores broadly."
        ],
        "Problem Addressed",
        [
            "Minimize the 2D Rastrigin function over the bounded domain [-5.12, 5.12].",
            "Compare three gradient descent variants against simulated annealing.",
            "Evaluate whether each method can reliably reach or approach the global minimum at f(0,0) = 0."
        ]
    )

    add_two_column_slide(
        prs,
        "Background And Approach",
        "Background",
        [
            "Rastrigin combines a quadratic bowl with cosine ripples, which creates many deceptive local minima.",
            "Gradient descent follows local slope information and can converge quickly when the step schedule is suitable.",
            "Simulated annealing can accept worse moves early, helping it escape local traps."
        ],
        "Methods Used",
        [
            "GD Fixed: constant learning rate alpha = 0.01.",
            "GD Decaying: alpha_t = alpha0 / (1 + decay * t) with alpha0 = 0.1 and decay = 0.001.",
            "GD Momentum: alpha = 0.01 and beta = 0.9.",
            "SA: T0 = 10, alpha = 0.995, step radius = 0.5."
        ]
    )

    add_bullets_slide(
        prs,
        "Experiment Setup",
        [
            "30 runs per algorithm from shared random starting points in the full domain.",
            "Maximum of 5000 iterations for every method.",
            "Success threshold: f(x) < 1e-3 counts as reaching the global minimum.",
            "Metrics: best objective value, success rate, distance to the optimum, and runtime.",
            "Using the same starting-point set made the comparison fair across algorithms."
        ]
    )

    add_table_slide(
        prs,
        "Evaluation Results",
        ["Algorithm", "Mean f", "Best f", "Success", "Time (ms)"],
        [
            ["GD_Fixed", rows["GD_Fixed"]["mean_f"], rows["GD_Fixed"]["best_f"], rows["GD_Fixed"]["success_rate"], rows["GD_Fixed"]["mean_time_ms"]],
            ["GD_Decaying", rows["GD_Decaying"]["mean_f"], rows["GD_Decaying"]["best_f"], rows["GD_Decaying"]["success_rate"], rows["GD_Decaying"]["mean_time_ms"]],
            ["GD_Momentum", rows["GD_Momentum"]["mean_f"], rows["GD_Momentum"]["best_f"], rows["GD_Momentum"]["success_rate"], rows["GD_Momentum"]["mean_time_ms"]],
            ["SA", rows["SA"]["mean_f"], rows["SA"]["best_f"], rows["SA"]["success_rate"], rows["SA"]["mean_time_ms"]],
        ],
        "Summary across 30 repeated trials"
    )

    add_image_slide(
        prs,
        "Landscape And Search Behavior",
        [project / "figures" / "rastrigin_contour.png", project / "figures" / "trajectories.png"],
        [
            "The contour plot shows the highly multimodal search landscape around the global minimum.",
            "Trajectory plots illustrate the difference between direct descent and exploratory stochastic search."
        ]
    )

    add_image_slide(
        prs,
        "Convergence And Distribution Of Results",
        [project / "figures" / "convergence_plot.png", project / "figures" / "boxplot.png"],
        [
            "Decaying GD converged consistently to near-zero values across all runs.",
            "The boxplot shows that fixed GD and momentum were less reliable, while SA often got close but rarely met the strict success threshold."
        ]
    )

    add_bullets_slide(
        prs,
        "Discussion, Conclusions, And Future Work",
        [
            "The best overall performer was decaying gradient descent, which achieved 30/30 successes with mean f approximately 1e-6.",
            "Simulated annealing improved exploration and found near-optimal points, but under the chosen schedule it did not satisfy the strict success threshold in any run.",
            "Fixed learning rate and momentum were faster than exhaustive exploration, but they remained vulnerable to local minima.",
            "Future work: tune the annealing schedule, test additional benchmark functions, and examine higher-dimensional Rastrigin settings."
        ]
    )

    out = project / "Rastrigin_project_presentation.pptx"
    prs.save(out)
    return out


def main():
    tools_dir = ROOT / "tools"
    tools_dir.mkdir(exist_ok=True)
    q_out = build_8queens()
    r_out = build_rastrigin()
    print(q_out)
    print(r_out)


if __name__ == "__main__":
    main()
